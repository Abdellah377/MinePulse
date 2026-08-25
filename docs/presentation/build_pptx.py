"""MinePulse presentation — AI / LangGraph first (screens are sensors, not the product)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(r"C:\Users\Gigabyte\Desktop\MinePulse")
SHOTS = ROOT / "docs" / "presentation" / "screenshots"
DIAG = ROOT / "docs" / "presentation" / "diagrams"
OUT = ROOT / "docs" / "presentation" / "MinePulse_Prototype.pptx"
LOGO = ROOT / "src" / "assets" / "ocp_logo.png"

GREEN = RGBColor(0x50, 0x80, 0x00)
GREEN_DARK = RGBColor(0x3D, 0x8C, 0x14)
INK = RGBColor(0x1C, 0x1D, 0x21)
MUTED = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF4, 0xF5, 0xF7)
CARD = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Inches(13.333), Inches(7.5)


def set_run(run, text, size=18, bold=False, color=INK):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    fill_shape(bg, BG)
    return slide


def add_bar(slide, color=GREEN):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.12))
    fill_shape(shape, color)


def heading(slide, text, top=0.35, size=26):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(0.55))
    set_run(box.text_frame.paragraphs[0].add_run(), text, size, True, INK)


def caption(slide, text, top=0.85):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(0.4))
    set_run(box.text_frame.paragraphs[0].add_run(), text, 13, False, MUTED)


def bullets(slide, lines, left=0.9, top=1.4, width=11.5, height=5.2, size=17):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left - 0.2), Inches(top - 0.25), Inches(width + 0.4), Inches(height)
    )
    fill_shape(card, CARD)
    try:
        card.adjustments[0] = 0.04
    except Exception:
        pass
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height - 0.3))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(11)
        set_run(p.add_run(), f"•  {line}", size, False, INK)


def fit_picture(slide, path: Path, left, top, max_w, max_h):
    if not path.exists():
        box = slide.shapes.add_textbox(left, top, max_w, Inches(0.5))
        set_run(box.text_frame.paragraphs[0].add_run(), f"[Image manquante: {path.name}]", 12, False, MUTED)
        return
    # add at max width; PowerPoint will keep aspect if we only set width
    pic = slide.shapes.add_picture(str(path), left, top, width=max_w)
    if pic.height > max_h:
        ratio = max_h / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = int(max_h)


def title_slide(prs):
    slide = blank_slide(prs)
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(3.4))
    fill_shape(banner, GREEN)
    if LOGO.exists():
        try:
            slide.shapes.add_picture(str(LOGO), Inches(0.7), Inches(0.45), height=Inches(0.65))
        except Exception:
            pass
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(12), Inches(1.8))
    tf = box.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(), "MinePulse AI", 46, True, WHITE)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    set_run(
        p2.add_run(),
        "LangGraph d’optimisation — anticipation & décisions pour plus de production",
        20,
        False,
        WHITE,
    )

    body = slide.shapes.add_textbox(Inches(0.7), Inches(3.8), Inches(12), Inches(2.8))
    tf = body.text_frame
    tf.word_wrap = True
    lines = [
        "Ce n’est pas un écran de supervision de plus.",
        "Les infos Film / Carte / Parc existent déjà chez l’opérateur.",
        "Le projet = un agent IA connecté à tout le poste, qui anticipe les problèmes",
        "et propose les meilleures décisions d’optimisation (humain au final).",
        "Démo : Khouribga — Merah El Ahrach · Poste matin · Mode prototype",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        set_run(p.add_run(), line, 16 if i else 18, i == 0, INK if i else GREEN_DARK)


def section(prs, title, subtitle=""):
    slide = blank_slide(prs)
    add_bar(slide)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.5), Inches(2))
    tf = box.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(), title, 38, True, GREEN_DARK)
    if subtitle:
        p = tf.add_paragraph()
        p.space_before = Pt(14)
        set_run(p.add_run(), subtitle, 18, False, MUTED)


def bullets_slide(prs, title, lines, sub=None):
    slide = blank_slide(prs)
    add_bar(slide)
    heading(slide, title)
    if sub:
        caption(slide, sub)
    bullets(slide, lines, top=1.35 if sub else 1.2)
    return slide


def diagram_slide(prs, title, sub, image: Path):
    slide = blank_slide(prs)
    add_bar(slide)
    heading(slide, title, size=24)
    caption(slide, sub)
    fit_picture(slide, image, Inches(0.7), Inches(1.25), Inches(12), Inches(5.8))


def screenshot_slide(prs, title, sub, image: Path):
    slide = blank_slide(prs)
    add_bar(slide)
    heading(slide, title, size=22)
    caption(slide, sub)
    fit_picture(slide, image, Inches(0.45), Inches(1.2), Inches(12.4), Inches(5.9))


def two_col(prs, title, left_t, left_b, right_t, right_b):
    slide = blank_slide(prs)
    add_bar(slide)
    heading(slide, title)
    for x, t, bs in ((0.6, left_t, left_b), (6.95, right_t, right_b)):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.15), Inches(5.8), Inches(5.7))
        fill_shape(card, CARD)
        try:
            card.adjustments[0] = 0.04
        except Exception:
            pass
        ht = slide.shapes.add_textbox(Inches(x + 0.3), Inches(1.4), Inches(5.2), Inches(0.5))
        set_run(ht.text_frame.paragraphs[0].add_run(), t, 17, True, GREEN_DARK)
        bx = slide.shapes.add_textbox(Inches(x + 0.3), Inches(2.0), Inches(5.2), Inches(4.6))
        tf = bx.text_frame
        tf.word_wrap = True
        for i, line in enumerate(bs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(10)
            set_run(p.add_run(), f"•  {line}", 14, False, INK)


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # --- AI thesis ---
    title_slide(prs)

    section(prs, "La vraie idée du projet", "Pas afficher des infos — les optimiser avec l’IA")

    bullets_slide(
        prs,
        "Ce que MinePulse n’est PAS",
        [
            "Un nouveau tableau Film / Carte / Parc pour remplacer OPM",
            "Une copie d’écrans que le chef de poste a déjà",
            "Un projet « dashboard » ou visualisation seule",
            "Une IA décorative (chips génériques sans décision)",
        ],
        sub="Les données opérationnelles existent déjà — on les réutilise comme entrée de l’agent",
    )

    bullets_slide(
        prs,
        "Ce que MinePulse EST",
        [
            "Un système IA type LangGraph : graphe d’agents connecté à tout le poste",
            "Il lit Film, Carte, Parc, Exceptions, Performance et Événements ensemble",
            "Il anticipe les problèmes futurs (congestion, attente, trou de production)",
            "Il propose les meilleures décisions d’optimisation pour récupérer du tonnage",
            "L’humain valide : Préparer / Marquer / Ignorer — jamais d’auto-application",
        ],
        sub="Produit = intelligence de décision · UI = capteurs + surfaces d’action",
    )

    diagram_slide(
        prs,
        "Sans AI vs Avec MinePulse AI",
        "De la réaction fragmentée → à l’anticipation + décision cohérente",
        DIAG / "ai-value-before-after.png",
    )

    section(prs, "Architecture IA", "LangGraph au centre — tous les écrans comme nœuds")

    diagram_slide(
        prs,
        "Orchestration LangGraph",
        "Un agent connecté à Film, Carte, Parc, Exceptions, Performance, Optimisation",
        DIAG / "ai-langgraph-architecture.png",
    )

    diagram_slide(
        prs,
        "Boucle de décision IA",
        "Télémétrie → Diagnostic → Anticipation → Options → Décision humaine",
        DIAG / "ai-decision-loop.png",
    )

    diagram_slide(
        prs,
        "Anticipation des problèmes",
        "Exemple : risque Banc B — agir avant que le trou de production s’aggrave",
        DIAG / "ai-anticipation-timeline.png",
    )

    diagram_slide(
        prs,
        "Humain dans la boucle",
        "L’IA propose — l’opérateur décide (Préparer / Marquer / Ignorer)",
        DIAG / "ai-human-in-loop.png",
    )

    section(prs, "Monde de démo", "Le scénario que l’agent raisonne — une seule vérité")

    two_col(
        prs,
        "Ce que l’agent voit (scénario Merah El Ahrach)",
        "Signaux",
        [
            "Production 7 231 / 8 160 t (−11 %)",
            "Banc B saturé : file 7 · cap. 3",
            "EXC-027 en maintenance (indisponible)",
            "TRK-012 arrêt non défini ~23 min",
            "TRK-004 perte télémétrie ~5 min",
            "Cycle moyen qui remonte après 11:00",
        ],
        "Décision IA prioritaire",
        [
            "Rediriger 3–4 camions Banc B → Banc A",
            "Ne jamais traiter EXC-027 comme « disponible »",
            "Estimer gain t/h si le plan est préparé",
            "Exposer confiance + preuves",
            "Laisser le chef de poste décider",
        ],
    )

    section(
        prs,
        "Les écrans = surfaces de l’IA",
        "Pas le produit — le graphe d’entrées / sorties de l’agent",
    )

    bullets_slide(
        prs,
        "Rôle de chaque surface",
        [
            "Exceptions — où l’IA explique « pourquoi » et pointe l’action",
            "Carte — contexte spatial (zones, files) pour le raisonnement",
            "Film — preuve temporelle des états (attente / arrêt)",
            "Parc — cycles et engins concernés par la reco",
            "Performance — mesure l’écart que l’IA cherche à combler",
            "Optimisation — cockpit de décision LangGraph (cœur produit)",
            "Événements — flux d’anomalies que l’agent corréle",
        ],
    )

    # Screenshots reframed as AI surfaces
    shots = [
        (
            "Surface IA — Optimisation (cœur)",
            "Situation → options → simulation → décision humaine (pas d’appliquer auto)",
            SHOTS / "mp-optimisation.png",
        ),
        (
            "Surface IA — Exceptions",
            "Brief + hypothèses / analyse IA sur sélection — pas un héros IA géant",
            SHOTS / "mp-exceptions.png",
        ),
        (
            "Entrée spatiale — Carte",
            "Carte MapLibre : l’agent localise congestion et engins spotlight",
            SHOTS / "mp-carte.png",
        ),
        (
            "Entrée temporelle — Film",
            "Preuve des états du poste pour le raisonnement (attente orange / arrêt rouge)",
            SHOTS / "mp-film.png",
        ),
        (
            "Entrée flotte — Parc",
            "Cycles et disponibilité : variables du modèle d’optimisation",
            SHOTS / "mp-parc.png",
        ),
        (
            "Mesure — Performance",
            "L’objectif que l’IA défend : 7 231 / 8 160 t, décrochage depuis 10:30",
            SHOTS / "mp-performance.png",
        ),
        (
            "Flux — Événements",
            "Anomalies live corrélées par l’agent (critique / warning / info)",
            SHOTS / "mp-evenements.png",
        ),
    ]
    for t, s, p in shots:
        screenshot_slide(prs, t, s, p)

    section(prs, "Roadmap IA", "Du prototype connecté → LangGraph production")

    two_col(
        prs,
        "Aujourd’hui vs demain",
        "Prototype actuel",
        [
            "Raisonnement scénario cohérent (mock)",
            "Bundle d’optimisation + slots IA contextuels",
            "Surfaces UI branchées sur une seule vérité",
            "Décision humaine explicite",
            "Carte / Film / Parc comme capteurs",
        ],
        "Cible LangGraph",
        [
            "Graphe multi-nœuds (diagnostic, anticipation, reco)",
            "Mémoire de poste + outils (zones, engins, events)",
            "Simulation what-if avant décision",
            "Connexion télémétrie / GIS OCP réelle",
            "Audit, confiance, et boucle d’apprentissage",
        ],
    )

    # Close
    slide = blank_slide(prs)
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.0), W, Inches(3.4))
    fill_shape(banner, GREEN)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.5), Inches(2.6))
    tf = box.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(), "MinePulse = AI d’optimisation", 36, True, WHITE)
    p = tf.add_paragraph()
    p.space_before = Pt(12)
    set_run(
        p.add_run(),
        "LangGraph connecté au poste · anticipe · décide mieux · plus de production\nLes écrans montrent — l’IA optimise.",
        18,
        False,
        WHITE,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
